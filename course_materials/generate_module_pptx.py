"""
Generates individual .pptx files for all 50 Melksham Mental Health course modules.
Run from the course_materials directory: python3 generate_module_pptx.py
Requires: pip install python-pptx Pillow
Background and logo images are extracted from extended_course_slides_with_user_logo.pptx
and saved as bg_opt.jpg and logo_opt.png before generation begins.
"""
import os

from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

ORANGE     = RGBColor(0xFF, 0x66, 0x00)
OFF_WHITE  = RGBColor(0xF0, 0xF0, 0xF0)
LIGHT_GREY = RGBColor(0xCC, 0xCC, 0xCC)
DARK_BG    = RGBColor(0x0A, 0x0A, 0x0A)

def add_text_box(slide, text, left, top, width, height, font_size, color,
                 bold=False, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    return txBox

def add_bullets(slide, bullets, left, top, width, height, font_size, color, title=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    if title:
        p = tf.paragraphs[0]; first = False
        r = p.add_run(); r.text = title
        r.font.size = Pt(font_size + 2)
        r.font.color.rgb = ORANGE; r.font.bold = True
    for b in bullets:
        p = tf.add_paragraph() if not first else tf.paragraphs[0]; first = False
        p.space_before = Pt(4)
        r = p.add_run(); r.text = "• " + b
        r.font.size = Pt(font_size); r.font.color.rgb = color
    return txBox

def set_bg(slide):
    bg = slide.background; fill = bg.fill
    fill.solid(); fill.fore_color.rgb = DARK_BG

def new_slide(prs, bg_img, logo_img):
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(sl)
    W = prs.slide_width.emu; H = prs.slide_height.emu
    M = Emu(457200)
    sl.shapes.add_picture(bg_img, 0, 0, W, H)
    logo_sz = Emu(1097280)
    sl.shapes.add_picture(logo_img, W - logo_sz - M, M // 2, logo_sz, logo_sz)
    return sl

def crisis_footer(s, W, H):
    add_text_box(s, "Crisis? Call 999 or Samaritans 116 123 (free, 24/7)",
                 Emu(457200), H - Emu(548640), W - Emu(914400), Emu(457200),
                 12, LIGHT_GREY, align=PP_ALIGN.CENTER)

def divider(s, W):
    d = s.shapes.add_shape(1, Emu(457200), Emu(1097280), W - Emu(914400), Emu(45720))
    d.fill.solid(); d.fill.fore_color.rgb = ORANGE; d.line.fill.background()

def build_pptx(num, title, bullets, refs, activity, bg_img, logo_img, out):
    prs = Presentation()
    prs.slide_width = Emu(9144000); prs.slide_height = Emu(6858000)
    W = prs.slide_width.emu; H = prs.slide_height.emu; M = Emu(457200)

    # Slide 1 - Title
    s = new_slide(prs, bg_img, logo_img)
    bar = s.shapes.add_shape(1, 0, 0, W, Emu(1828800))
    bar.fill.solid(); bar.fill.fore_color.rgb = ORANGE; bar.line.fill.background()
    add_text_box(s, f"Module {num}", M, Emu(274320), W - Emu(1554480), Emu(548640), 32, DARK_BG, bold=True)
    add_text_box(s, title, M, Emu(640080), W - Emu(1554480), Emu(914400), 26, DARK_BG)
    add_text_box(s, "Melksham Mental Health — Comprehensive Course",
                 M, Emu(2057400), W - 2*M, Emu(548640), 20, ORANGE, bold=True)
    add_text_box(s, "2-hour session  •  Includes 30-minute group discussion",
                 M, Emu(2606040), W - 2*M, Emu(457200), 16, LIGHT_GREY)
    crisis_footer(s, W, H)

    # Slide 2 - Session Overview
    s = new_slide(prs, bg_img, logo_img)
    add_text_box(s, "Session Overview", M, M, W - 2*M, Emu(640080), 28, ORANGE, bold=True)
    divider(s, W)
    add_bullets(s, [
        "Introduction and Objectives  (15 min) — Key terms and learning goals",
        "Evidence-Based Content  (45 min) — Research, risk factors and treatments",
        "Interactive Activity  (30 min) — Polls, reflection and role-play",
        "Group Discussion  (30 min) — Confidential sharing and peer learning",
        "Wrap-Up and Resources  (15 min) — Summary, helplines and further reading",
    ], M, Emu(1280160), W - 2*M, Emu(4572000), 18, OFF_WHITE)
    crisis_footer(s, W, H)

    # Slide 3 - Learning Objectives
    s = new_slide(prs, bg_img, logo_img)
    add_text_box(s, "Learning Objectives", M, M, W - 2*M, Emu(640080), 28, ORANGE, bold=True)
    divider(s, W)
    add_bullets(s, ["Understand " + (b if len(b) <= 80 else b[:b.rfind(' ', 0, 80)] or b[:80]) for b in bullets] + [
        "Apply coping strategies and evidence-based techniques",
        "Identify when and how to access professional support",
    ], M, Emu(1280160), W - 2*M, Emu(4572000), 18, OFF_WHITE)
    crisis_footer(s, W, H)

    # Slide 4 - Key Topics
    s = new_slide(prs, bg_img, logo_img)
    add_text_box(s, f"Module {num}: Key Topics", M, M, W - 2*M, Emu(640080), 28, ORANGE, bold=True)
    divider(s, W)
    add_bullets(s, bullets, M, Emu(1280160), W - 2*M, Emu(4572000), 18, OFF_WHITE)
    crisis_footer(s, W, H)

    # Slide 5 - Evidence
    s = new_slide(prs, bg_img, logo_img)
    add_text_box(s, "Evidence and Research", M, M, W - 2*M, Emu(640080), 28, ORANGE, bold=True)
    divider(s, W)
    add_bullets(s, [
        "Topic: " + title,
        "Evidence-based approaches are used throughout this session",
        "Content drawn from NHS guidelines, NICE, WHO and peer-reviewed research",
        "UK-specific statistics and resources are prioritised",
        "All content reviewed for accuracy and safety",
    ], M, Emu(1280160), W - 2*M, Emu(3657600), 18, OFF_WHITE)
    add_text_box(s, "This content does not replace professional mental health advice.",
                 M, Emu(5486400), W - 2*M, Emu(457200), 14, ORANGE, bold=True)
    crisis_footer(s, W, H)

    # Slide 6 - Interactive Activity
    s = new_slide(prs, bg_img, logo_img)
    add_text_box(s, "Interactive Activity", M, M, W - 2*M, Emu(640080), 28, ORANGE, bold=True)
    divider(s, W)
    add_text_box(s, activity, M, Emu(1463040), W - 2*M, Emu(2743200), 20, OFF_WHITE, wrap=True)
    add_text_box(s, "Instructions:", M, Emu(4114800), W - 2*M, Emu(365760), 16, ORANGE, bold=True)
    add_bullets(s, [
        "Take 5 minutes to reflect individually",
        "Share your thoughts with the group (optional — sharing is never required)",
        "Facilitator will guide the discussion",
    ], M, Emu(4480560), W - 2*M, Emu(1371600), 16, LIGHT_GREY)
    crisis_footer(s, W, H)

    # Slide 7 - Group Discussion
    s = new_slide(prs, bg_img, logo_img)
    add_text_box(s, "Group Discussion (30 minutes)", M, M, W - 2*M, Emu(640080), 28, ORANGE, bold=True)
    divider(s, W)
    add_bullets(s, [
        "What stood out most to you about " + title.lower() + "?",
        "Has this topic affected you, or someone you know?",
        "What barriers prevent people from seeking help with this?",
        "What one thing could you do differently after today?",
        "What support is available in our community?",
    ], M, Emu(1280160), W - 2*M, Emu(3657600), 18, OFF_WHITE,
    title="Reflection Questions")
    add_text_box(s, "Remember: You are never alone. Help is always available.",
                 M, Emu(5120640), W - 2*M, Emu(457200), 14, ORANGE, bold=True, align=PP_ALIGN.CENTER)
    crisis_footer(s, W, H)

    # Slide 8 - Resources
    s = new_slide(prs, bg_img, logo_img)
    add_text_box(s, "Resources and Further Support", M, M, W - 2*M, Emu(640080), 28, ORANGE, bold=True)
    divider(s, W)
    add_bullets(s, refs, M, Emu(1280160), W - 2*M, Emu(2743200), 18, OFF_WHITE, title="Key References")
    add_bullets(s, [
        "Emergency: 999",
        "NHS 111: 24/7 urgent mental health support",
        "Samaritans: 116 123 (free, 24/7)",
        "Shout Crisis Text: Text SHOUT to 85258",
        "Wiltshire Crisis Line: 0800 953 0110 (24/7)",
    ], M, Emu(4297200), W - 2*M, Emu(1828800), 14, LIGHT_GREY, title="Crisis Support")
    crisis_footer(s, W, H)

    # Slide 9 - Wrap-Up
    s = new_slide(prs, bg_img, logo_img)
    bar = s.shapes.add_shape(1, 0, 0, W, Emu(1828800))
    bar.fill.solid(); bar.fill.fore_color.rgb = ORANGE; bar.line.fill.background()
    add_text_box(s, f"Module {num} Complete", M, Emu(274320), W - Emu(1554480), Emu(640080), 30, DARK_BG, bold=True)
    add_text_box(s, title, M, Emu(731520), W - Emu(1554480), Emu(548640), 22, DARK_BG)
    add_bullets(s, [
        f"You have completed Module {num}: " + title,
        "Key concepts and evidence-based strategies have been explored",
        "Please complete the post-session reflection form",
        "Access course materials and handouts in the Members Portal",
        "Our next session builds on today's learning",
    ], M, Emu(2057400), W - 2*M, Emu(3200400), 18, OFF_WHITE)
    add_text_box(s, "melksham-mentalhealth.us  |  Melksham-mental-health@outlook.com",
                 M, H - Emu(640080), W - 2*M, Emu(457200), 12, ORANGE, align=PP_ALIGN.CENTER)

    prs.save(out)

MODULES = [
    (1, "Foundations of Mental Health",
     ["Define mental health and mental illness",
      "Discuss prevalence (970 million people with mental disorders) and economic burden",
      "Examine stigma and importance of self-care and social connections",
      "Explore WHO reports on workplace mental health risk"],
     ["WHO World Mental Health Report 2022", "NHS Mental Health Statistics", "Mind UK: What is mental health?"],
     "Reflection exercise: What does 'mental health' mean to you personally?"),
    (2, "Introduction to Mental-Health Disorders",
     ["Overview of DSM-5 categories: mood, anxiety, trauma-related, neurodevelopmental",
      "Neurocognitive, psychotic, personality, somatic, dissociative, eating, sleep disorders",
      "Impulse-control and substance-use disorders overview",
      "Preparing participants for deeper dives in later modules"],
     ["DSM-5 Diagnostic and Statistical Manual", "WHO ICD-11 Classification", "MIND UK: Types of mental health problems"],
     "Quiz: Match common symptoms to disorder categories"),
    (3, "Neurodevelopmental Disorders: Autism Spectrum",
     ["Autism as a spectrum of differences in brain development",
      "Early signs: limited eye contact, repetitive motions, sensory sensitivities",
      "Therapies to support communication and social skills",
      "Inclusive language and autistic-affirming practices"],
     ["WHO on Autism Spectrum Disorder", "National Autistic Society (NAS) UK", "Cleveland Clinic: Autism overview"],
     "Group exercise: Empathy mapping and inclusive language practice"),
    (4, "Neurodevelopmental Disorders: ADHD",
     ["Attention-deficit hyperactivity disorder explained",
      "Common symptoms: daydreaming, forgetfulness, hyperactivity, impulsive decisions",
      "Treatment strategies: medication, behaviour therapy, coaching",
      "Coping techniques for adults with ADHD"],
     ["CDC: ADHD Symptoms and Diagnosis", "ADHD UK: Living with ADHD", "NHS: ADHD Treatment"],
     "Interactive: Personal strategy mapping for attention and focus"),
    (5, "Neurodevelopmental Disorders: Intellectual and Learning Disabilities",
     ["Intellectual disability and global developmental delay",
      "Dyslexia and other specific learning disorders",
      "Early intervention and inclusive education strategies",
      "Supporting adults with learning disabilities"],
     ["VeryWellMind: Intellectual Developmental Disorder", "British Dyslexia Association", "Mencap UK Resources"],
     "Scenario discussion: Creating inclusive learning environments"),
    (6, "Trauma and Stress: Acute Stress and Adjustment Disorders",
     ["Defining acute stress disorder and adjustment disorders",
      "Major life changes (job loss, divorce) triggering anxiety and depressed mood",
      "Coping strategies and trauma-informed listening techniques",
      "When to seek professional support"],
     ["VeryWellMind: Acute Stress Disorder", "APA: Adjustment Disorders", "Mind UK: Stress"],
     "Role-play: Trauma-informed conversations"),
    (7, "Trauma and Stress: PTSD and Complex PTSD",
     ["PTSD symptoms: re-experiencing, avoidance, hypervigilance",
      "Complex PTSD (CPTSD) from prolonged or repeated trauma",
      "Evidence-based treatments: cognitive processing therapy, prolonged exposure",
      "Trauma-focused CBT and EMDR overview"],
     ["APA PTSD Treatment Guidelines", "NICE Guidelines: PTSD", "Combat Stress UK: Veterans PTSD"],
     "Grounding techniques practice: 5-4-3-2-1 sensory method"),
    (8, "Trauma and Stress: Domestic and Intimate Partner Violence",
     ["Prevalence: over one-third of women and one-quarter of men affected",
      "Types of abuse: physical, sexual, psychological, financial, coercive control",
      "Risk factors: young age, low income, social isolation",
      "Mental-health consequences: PTSD, depression, anxiety"],
     ["SafeLives UK: Domestic Abuse Statistics", "Womens Aid UK Resources", "NHS: Domestic Violence Support"],
     "Safety planning resources and referral pathways"),
    (9, "Trauma and Stress: Childhood Trauma and ACEs",
     ["Adverse childhood experiences (ACEs): abuse, neglect, household dysfunction",
      "How early trauma affects brain development",
      "Increased risk of mental disorders and chronic disease",
      "Resilience, protective factors and post-traumatic growth"],
     ["CDC: ACEs Study", "NSPCC: Child Abuse and Mental Health", "YoungMinds UK: Trauma in Children"],
     "Resilience-building exercise: Identifying protective factors"),
    (10, "Trauma and Stress: Self-harm and Suicide Prevention",
     ["Understanding reasons for self-harm: relief of distress, self-punishment",
      "Common behaviours and harm-reduction strategies",
      "Recognising warning signs in others",
      "Safety planning and crisis resources"],
     ["Samaritans: 116 123 (24/7)", "Crisis Text Line: Text SHOUT to 85258", "PAPYRUS UK: Young Peoples Suicide Prevention"],
     "Safety planning exercise and practising supportive conversations"),
    (11, "Mood Disorders: Major Depression",
     ["Symptoms: sad mood, loss of interest, fatigue, guilt, suicidal thoughts",
      "Causes: biological, psychological, social factors",
      "Treatments: antidepressants, psychotherapy, lifestyle changes",
      "Depression across cultures and genders"],
     ["NIMH: Depression Facts", "NHS: Clinical Depression", "Mind UK: Depression"],
     "Mood tracking and behavioural activation techniques"),
    (12, "Mood Disorders: Bipolar Disorder",
     ["Bipolar I and II disorders: key differences",
      "Manic episodes: elevated mood, reduced sleep, impulsivity",
      "Depressive episodes: low mood, fatigue, hopelessness",
      "Medication adherence and early warning signs"],
     ["VeryWellMind: Bipolar Disorder", "Bipolar UK: Support Resources", "NHS: Bipolar Disorder Treatment"],
     "Mood journaling and episode pattern recognition"),
    (13, "Mood Disorders: Seasonal and Persistent Depressive Disorders",
     ["Seasonal affective disorder (SAD): winter-pattern depression",
      "Dysthymia (persistent depressive disorder): chronic low mood",
      "Cyclothymia: mild mood fluctuations",
      "Light therapy, behavioural activation, and seasonal coping"],
     ["NHS: Seasonal Affective Disorder", "Mind UK: Seasonal Depression", "SADA: Seasonal Affective Disorder Association"],
     "Seasonal wellness planning: building a winter wellbeing toolkit"),
    (14, "Anxiety Disorders: Generalised Anxiety Disorder",
     ["Excessive worry, restlessness, concentration problems, muscle tension",
      "Physical symptoms: headaches, fatigue, insomnia",
      "Cognitive behavioural therapy (CBT) for GAD",
      "Mindfulness, medication, and lifestyle approaches"],
     ["NICE Guidelines: Generalised Anxiety Disorder", "Anxiety UK: GAD Support", "NHS: Generalised Anxiety Disorder"],
     "Worry time technique and cognitive restructuring exercise"),
    (15, "Anxiety Disorders: Phobias and Panic Disorder",
     ["Specific phobias: triggers, avoidance and impact on daily life",
      "Social anxiety: fear of judgment and embarrassment",
      "Panic disorder: unexpected panic attacks and anticipatory anxiety",
      "Exposure therapy, relaxation exercises and breathing techniques"],
     ["VeryWellMind: Phobias", "NHS: Panic Disorder", "No Panic UK: Support Helpline"],
     "Graduated exposure hierarchy: creating a fear ladder"),
    (16, "Anxiety Disorders: OCD and Related Conditions",
     ["Obsessive-compulsive disorder: intrusive thoughts and compulsive behaviours",
      "Body dysmorphic disorder, hoarding disorder, trichotillomania",
      "Exposure and response prevention (ERP) therapy",
      "Living with OCD: reducing shame and stigma"],
     ["OCD UK: Support Resources", "VeryWellMind: OCD-Related Disorders", "NHS: OCD Treatment"],
     "ERP practice: identifying triggers and resisting compulsions"),
    (17, "Dissociative and Somatic Disorders",
     ["Dissociative amnesia, dissociative identity disorder, depersonalisation",
      "Somatic symptom disorder and illness anxiety",
      "Connection between trauma and dissociation",
      "Trauma-focused approaches and grounding techniques"],
     ["ISSTD: Dissociative Disorders", "VeryWellMind: Somatic Symptom Disorder", "NHS: Dissociative Disorders"],
     "Grounding techniques: body-based and sensory anchoring exercises"),
    (18, "Eating Disorders and Body Image",
     ["Anorexia nervosa, bulimia nervosa, binge eating disorder",
      "Avoidant restrictive food intake disorder (ARFID)",
      "Body image and social media influences",
      "Recovery options and specialist support"],
     ["Beat Eating Disorders UK: 0808 801 0677", "NHS: Eating Disorders", "Mind UK: Eating Problems"],
     "Media literacy exercise: challenging unrealistic body standards"),
    (19, "Sleep and Circadian Rhythm Disorders",
     ["Insomnia, hypersomnolence, narcolepsy, sleep apnea",
      "Parasomnias and restless legs syndrome",
      "Sleep hygiene fundamentals and CBT for insomnia (CBT-I)",
      "Circadian rhythm disruption and mental health impacts"],
     ["Sleep Foundation: Sleep Disorders", "NHS: Insomnia", "NICE Guidelines: CBT for Insomnia"],
     "Sleep diary exercise and sleep hygiene plan creation"),
    (20, "Impulse-Control and Disruptive Disorders",
     ["Intermittent explosive disorder: anger outbursts",
      "Conduct disorder and oppositional defiant disorder in youth",
      "Kleptomania and pyromania: understanding impulse dysregulation",
      "Behavioural interventions and family therapy approaches"],
     ["VeryWellMind: Impulse Control Disorders", "NHS: Conduct Disorder", "YoungMinds: Behavioural Problems"],
     "Anger management techniques and impulse control strategies"),
    (21, "Substance Use and Addiction",
     ["Substance use disorder criteria and brain changes",
      "Alcohol, opioids, stimulants, cannabis and gambling disorders",
      "Harm-reduction strategies and overdose prevention",
      "Recovery pathways: 12-step, SMART Recovery, medication-assisted"],
     ["We Are With You UK", "Alcoholics Anonymous UK: 0800 9177 650", "Frank UK: Drug Information"],
     "Recovery planning: identifying triggers, supports and coping strategies"),
    (22, "Psychotic Disorders: Schizophrenia and Related Conditions",
     ["Positive symptoms: hallucinations, delusions, disorganised thinking",
      "Negative symptoms: flat affect, poverty of speech, avolition",
      "Early intervention in psychosis (EIP) services",
      "Antipsychotic medications and psychosocial support"],
     ["Rethink Mental Illness UK", "NHS Early Intervention in Psychosis", "Mind UK: Schizophrenia"],
     "Psychoeducation: early warning signs and relapse prevention planning"),
    (23, "Personality Disorders: Cluster A and B",
     ["Cluster A: paranoid, schizoid, schizotypal personality disorders",
      "Cluster B: antisocial, borderline, histrionic, narcissistic PD",
      "Dialectical behaviour therapy (DBT) for borderline PD",
      "Mentalisation-based treatment and reducing stigma"],
     ["Mind UK: Personality Disorders", "Emergence UK: BPD Support", "NICE Guidelines: Borderline Personality Disorder"],
     "DBT skills practice: distress tolerance and emotion regulation"),
    (24, "Personality Disorders: Cluster C and Others",
     ["Cluster C: avoidant, dependent, obsessive-compulsive PD",
      "Narcissistic personality disorder: recognition and boundaries",
      "Cognitive and interpersonal therapy approaches",
      "Self-compassion and healthy relationship building"],
     ["Mind UK: Personality Disorders", "VeryWellMind: Cluster C Personality Disorders", "NHS: Personality Disorder Treatment"],
     "Schema identification exercise: recognising maladaptive patterns"),
    (25, "Neurocognitive Disorders",
     ["Delirium: acute confusion, causes, recognition, management",
      "Major neurocognitive disorder (dementia) due to Alzheimers disease",
      "Parkinsons disease dementia and Lewy body dementia",
      "Caregiver support and advance care planning"],
     ["Alzheimers Society UK: 0333 150 3456", "Dementia UK: Admiral Nurse Service", "NHS: Dementia"],
     "Caregiver wellbeing: identifying support needs and self-care strategies"),
    (26, "Mental Health and Chronic Physical Conditions",
     ["Chronic illness (diabetes, arthritis, cancer) and mental health",
      "People with long-term conditions twice as likely to develop depression",
      "Bidirectional relationship: mental health affects physical recovery",
      "Coping strategies, peer support groups and self-management"],
     ["NHS: Long-term Conditions and Mental Health", "Mind UK: Physical Health and Mental Health", "Macmillan Cancer Support"],
     "Coping skills workshop: adapting to life with chronic illness"),
    (27, "Workplace Mental Health and Burnout",
     ["Poor working environments: discrimination, overwork, low control",
      "Job insecurity and mental health risks",
      "Burnout: exhaustion, cynicism, reduced efficacy",
      "Employer and employee strategies for prevention"],
     ["WHO: Mental Health in the Workplace", "HSE: Work-related Stress", "Mind UK: Workplace Wellbeing"],
     "Stress audit: mapping workplace stressors and identifying solutions"),
    (28, "Social Isolation and Loneliness",
     ["Difference between loneliness and social isolation",
      "Loneliness doubles risk of depression; increases anxiety and suicide risk",
      "Populations at greatest risk: older adults, young people, carers",
      "Interventions to strengthen social connections"],
     ["Campaign to End Loneliness UK", "Mind UK: Loneliness", "NHS: Loneliness and Mental Health"],
     "Social mapping: identifying connection points and community resources"),
    (29, "Financial Stress and Mental Health",
     ["Mental health issues leading to overspending and bill avoidance",
      "Money problems triggering anxiety, panic and shame",
      "The debt-mental health cycle",
      "Budgeting tools, benefits advice and debt support"],
     ["Mind UK: Money and Mental Health", "Money and Mental Health Policy Institute", "Citizens Advice UK: Debt Help"],
     "Financial wellbeing planning: small steps toward financial stability"),
    (30, "Relationships and Attachment",
     ["Healthy relationship skills: communication, empathy, boundaries",
      "Effects of codependency and insecure attachment styles",
      "How trauma and mental illness affect partnerships",
      "Building secure and supportive relationships"],
     ["Relate UK: Relationship Support", "Mind UK: Relationships and Mental Health", "NHS: Relationship Counselling"],
     "Attachment style reflection and healthy communication practice"),
    (31, "Parenting Stress and Paternal Mental Health",
     ["Mental-health impacts of parenting: sleep deprivation, identity shift",
      "One in ten dads experience postpartum depression",
      "5 to 15 percent of new fathers develop anxiety disorders",
      "Coping strategies and support for parents"],
     ["Pandas Foundation UK: Perinatal Mental Health", "Dad Matters UK", "NHS: Mental Health for New Dads"],
     "Peer support circles: normalising parenting struggles"),
    (32, "Maternal Mental Health and Perinatal Disorders",
     ["Postpartum depression: symptoms, prevalence, treatment",
      "Perinatal anxiety, OCD and postpartum psychosis",
      "Diagnosis rates doubled between 2010 and 2021",
      "Screening, early intervention and specialist support"],
     ["PANDAS Foundation UK: 0808 1961 776", "APNI: Association for Post Natal Illness", "NHS: Perinatal Mental Health"],
     "Wellbeing planning for new mothers: building support networks"),
    (33, "Child and Adolescent Mental Health",
     ["Common issues: anxiety, depression, self-harm, behavioural disorders",
      "Risk factors: bullying, academic pressure, social media",
      "Early intervention and CAMHS services",
      "Supporting young peoples mental health at home and school"],
     ["YoungMinds UK: 0808 802 5544", "CAMHS UK Resources", "Place2Be: School Mental Health"],
     "Workshop: How to talk to young people about mental health"),
    (34, "Older Adult Mental Health",
     ["14.1 percent of adults aged 70 and over have a mental disorder",
      "One sixth of suicide deaths occur in older adults",
      "Risk factors: bereavement, declining health, social isolation",
      "Elder abuse and its mental health consequences"],
     ["Age UK: 0800 678 1602", "Mental Health Foundation: Older People", "NHS: Mental Health in Later Life"],
     "Life review therapy: reminiscence and meaning-making exercises"),
    (35, "LGBTQIA Plus Mental Health",
     ["Higher rates of depression, self-harm and suicidal thoughts in LGBTQIA+ people",
      "Caused by discrimination, homophobia, transphobia and rejection",
      "Minority stress model: chronic stressors from marginalisation",
      "Affirmative therapy and community support"],
     ["Stonewall UK: LGBTQ+ Mental Health", "Mermaids UK: Trans Young People", "LGBT Foundation: 0345 3 30 30 30"],
     "Allyship training: affirmative language and supportive environments"),
    (36, "Racism, Discrimination and Intersectionality",
     ["Racism as a mental-health issue: trauma and ill-health",
      "Compounded effects of racism, sexism, ableism and discrimination",
      "Intersectionality: multiple marginalised identities",
      "Community resilience and culturally-sensitive support"],
     ["Rethink Mental Illness: Racism and Mental Health", "BAME Mental Health UK", "NHS Race and Health Observatory"],
     "Reflective dialogue: exploring privilege, bias and solidarity"),
    (37, "Gender and Sexuality",
     ["Gender identity and expectations influencing mental health",
      "Gender dysphoria: diagnosis, treatment and affirming care",
      "Sexism, gender-based violence and their psychological impact",
      "Inclusive language and allyship training"],
     ["Gendered Intelligence UK", "Gender Identity Clinic NHS", "Stonewall UK: Trans Mental Health"],
     "Gender-affirming conversation practice and inclusive language guide"),
    (38, "Sexual Violence and Trauma",
     ["Sexual assault survivors often develop PTSD, anxiety and depression",
      "Hypervigilance, flashbacks, nightmares and self-blame",
      "Emotional dysregulation and relationship difficulties",
      "Trauma-focused therapies and specialist support services"],
     ["Rape Crisis England and Wales: 0808 500 2222", "RASASC: Rape and Sexual Abuse Support Centre", "Mind UK: Sexual Abuse and Mental Health"],
     "Trauma-informed support: how to respond to disclosures safely"),
    (39, "Grief, Loss and Bereavement",
     ["Stages of grief: denial, anger, bargaining, depression, acceptance",
      "Complicated grief and prolonged grief disorder",
      "Ambiguous loss: missing persons, estrangement, dementia",
      "Disenfranchised grief: losses not socially recognised"],
     ["Cruse Bereavement Care UK: 0808 808 1677", "Child Bereavement UK", "Mind UK: Bereavement"],
     "Grief ritual activity: creating a memory or tribute"),
    (40, "Chronic Pain and Pain Management",
     ["Interplay between chronic pain and mental health",
      "Increased risk of depression and anxiety with chronic pain",
      "Cognitive behavioural therapy (CBT) for chronic pain",
      "Pacing, mindfulness and acceptance-based approaches"],
     ["Pain UK: Chronic Pain Resources", "British Pain Society", "NHS: Chronic Pain Management"],
     "Pain management planning: activity pacing and coping strategies"),
    (41, "Sleep Hygiene and Circadian Health",
     ["Strategies for healthy sleep patterns and bedtime routines",
      "Impact of shift work and irregular schedules on circadian rhythms",
      "Screen time and blue light disrupting melatonin production",
      "Sleep restriction therapy and relaxation exercises for insomnia"],
     ["Sleep Foundation UK", "NHS: Sleep and Tiredness", "NICE: Insomnia Management Guidelines"],
     "Sleep hygiene audit: creating a personalised sleep improvement plan"),
    (42, "Nutrition, Exercise and Mental Health",
     ["Regular physical activity improves mood and cognitive function",
      "Balanced diet and hydration for mental wellbeing",
      "Nutritional deficiencies: omega-3, vitamin D and B vitamins",
      "The gut-brain axis and microbiome role in mental health"],
     ["Mental Health Foundation: Food and Mood", "NHS: Exercise for Depression", "Mind UK: Physical Activity and Mental Health"],
     "Food and mood diary: identifying dietary patterns affecting wellbeing"),
    (43, "Mindfulness, Meditation and Breathwork",
     ["Mindfulness-based stress reduction (MBSR) principles",
      "Breathing techniques: diaphragmatic, box breathing, 4-7-8",
      "Yoga, body scan and meditation practices",
      "Evidence for reducing anxiety and enhancing emotional regulation"],
     ["Mindfulness-Based Cognitive Therapy (MBCT) - NICE approved", "Breathworks UK: Mindfulness", "Oxford Mindfulness Centre"],
     "Guided mindfulness practice: 10-minute body scan and breath focus"),
    (44, "Creative Arts and Expressive Therapies",
     ["Art therapy: using visual art to process emotions and trauma",
      "Music therapy: rhythm, sound and emotional expression",
      "Dance and movement therapy: body-based processing",
      "Writing and journalling as therapeutic tools"],
     ["British Association of Art Therapists", "British Association for Music Therapy", "Lapidus UK: Words and Wellbeing"],
     "Creative expression workshop: participants create a visual or written piece"),
    (45, "Technology and Mental Health",
     ["Effects of excessive smartphone use on anxiety, depression and sleep",
      "Gaming disorder: criteria, prevalence and treatment",
      "Social media comparison, cyberbullying and body image",
      "Digital detox strategies and healthy online engagement"],
     ["NHS: Social Media and Young Peoples Mental Health", "Cybersmile Foundation", "Centre for Humane Technology"],
     "Digital wellness audit: tracking screen time and setting healthy boundaries"),
    (46, "Climate Anxiety and Eco-distress",
     ["Mental-health impacts of climate change: fear, grief, helplessness",
      "Climate anxiety and eco-distress across age groups",
      "Eco-therapy and nature-based interventions",
      "Activism, community resilience and meaningful action"],
     ["Climate Psychology Alliance UK", "Good Grief Network", "Mind UK: Nature and Mental Health"],
     "Group discussion: transforming climate grief into purposeful action"),
    (47, "Spirituality and Meaning",
     ["Spirituality, religion and personal meaning supporting mental health",
      "Research on spiritual wellbeing and psychological resilience",
      "Spiritual abuse and faith-related trauma",
      "Supporting clients with diverse spiritual and religious beliefs"],
     ["Royal College of Psychiatrists: Spirituality", "Spirituality and Mental Health Forum", "Samaritans: 116 123"],
     "Values clarification exercise: exploring sources of meaning and purpose"),
    (48, "Mens Mental Health and Masculinity",
     ["Traditional masculinity norms hindering help-seeking",
      "Mens mental health statistics: higher suicide rates, less treatment",
      "Anger, substance use and risk-taking as hidden distress signals",
      "Gender-sensitive outreach and peer support strategies"],
     ["Movember Foundation UK", "Andys Man Club UK", "CALM: Campaign Against Living Miserably 0800 585858"],
     "Mens peer group discussion: redefining strength and asking for help"),
    (49, "Military and Veteran Mental Health",
     ["Combat-related PTSD: prevalence, symptoms and treatment",
      "Moral injury: guilt and shame from combat experiences",
      "Traumatic brain injury (TBI) and its mental health consequences",
      "Reintegration challenges and substance use in veterans"],
     ["Combat Stress UK: 0800 138 1619", "Veterans UK: 0808 1914 218", "Help for Heroes UK"],
     "Reintegration planning: building civilian support networks"),
    (50, "Prison, Legal System and Mental Health",
     ["Mental-health challenges faced by incarcerated individuals",
      "Impact of solitary confinement on psychological wellbeing",
      "Mental health of those on probation and with criminal records",
      "Stigma, reentry support and community rehabilitation"],
     ["Prison Reform Trust UK", "Centre for Mental Health: Criminal Justice", "NACRO: Crime Reduction Charity"],
     "Reentry support planning: identifying mental health resources post-release"),
]

BG   = os.path.join(os.path.dirname(__file__), "bg_opt.jpg")
LOGO = os.path.join(os.path.dirname(__file__), "logo_opt.png")
OUT  = os.path.join(os.path.dirname(__file__), "modules")
os.makedirs(OUT, exist_ok=True)

print(f"Generating {len(MODULES)} PowerPoint files...")
for n, title, bullets, refs, activity in MODULES:
    # Sanitise title for use as a filename component; keep only word chars, spaces, hyphens
    import re
    safe = re.sub(r'[^\w\s-]', '', title).strip().replace(' ', '_')
    # Limit to 60 chars to stay well within filesystem limits while avoiding mid-word cuts
    if len(safe) > 60:
        safe = safe[:safe.rfind('_', 0, 60)] or safe[:60]
    fname = f"Module_{n:02d}_{safe}.pptx"
    path  = os.path.join(OUT, fname)
    build_pptx(n, title, bullets, refs, activity, BG, LOGO, path)
    print(f"  [{n:2d}/50] {fname}")

print(f"\nAll {len(MODULES)} files saved to {OUT}")
